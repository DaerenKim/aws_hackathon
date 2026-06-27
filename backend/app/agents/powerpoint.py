"""PowerPoint Agent for Hackathon Studio.

Generates a hackathon presentation using python-pptx with 6-15 slides
covering: problem statement, solution overview, technical architecture,
demo screenshots, team/tooling, and future roadmap. Each slide includes
speaker notes of at least 50 characters.

Requirements: 11.1, 11.2, 11.3, 11.4, 11.5, 11.6, 11.7
"""

import io
import json
import logging
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.agents.base import BaseAgent
from app.models.artifacts import AgentResult

logger = logging.getLogger(__name__)

SYSTEM_PROMPT = (
    "You are a presentation designer. Create compelling hackathon pitch slides. "
    "Your output must be structured JSON that can be directly parsed."
)

# Minimum required slide topics per Requirement 11.2
REQUIRED_TOPICS = [
    "problem_statement",
    "solution_overview",
    "technical_architecture",
    "demo_screenshots",
    "team_tooling",
    "future_roadmap",
]

MIN_SLIDES = 6
MAX_SLIDES = 15
MIN_SPEAKER_NOTES_LENGTH = 50

OUTPUT_DIR = "ppt/"
PPTX_FILENAME = "ppt/presentation.pptx"
SPEAKER_NOTES_FILENAME = "ppt/speaker_notes.md"


class PowerPointAgent(BaseAgent):
    """Agent that generates a hackathon presentation using python-pptx.

    Reads project_spec.md and architecture.md from the shared workspace,
    uses the LLM to generate slide content (titles, bullet points, speaker
    notes), then creates a .pptx file with 6-15 slides.

    Requirements: 11.1, 11.2, 11.3, 11.4, 11.5, 11.6, 11.7
    """

    async def execute(self, context: dict) -> AgentResult:
        """Execute presentation generation.

        Args:
            context: Dictionary with optional overrides. Currently unused.

        Returns:
            AgentResult indicating success/failure with artifacts produced.
        """
        start_time = time.monotonic()
        artifacts_produced: list[str] = []
        missing_artifacts: list[str] = []

        # Step 1: Read source artifacts from workspace
        project_spec = await self._read_source_artifact(
            "project_spec.md", missing_artifacts
        )
        architecture = await self._read_source_artifact(
            "architecture.md", missing_artifacts
        )

        # Read optional artifacts for richer content
        judge_analysis = await self._read_source_artifact(
            "judge_analysis.md", missing_artifacts, required=False
        )
        roadmap = await self._read_source_artifact(
            "roadmap.md", missing_artifacts, required=False
        )

        # Requirement 11.6: report missing required artifacts
        if "project_spec.md" in missing_artifacts and "architecture.md" in missing_artifacts:
            await self.log(
                "Both project_spec.md and architecture.md are missing. "
                "Generating presentation with minimal placeholder content."
            )

        # Step 2: Check for demo screenshots
        screenshots = await self._find_screenshots()

        # Step 3: Generate slide content via LLM
        await self.log("Generating slide content via LLM...")
        slide_content = await self._generate_slide_content(
            project_spec=project_spec,
            architecture=architecture,
            judge_analysis=judge_analysis,
            roadmap=roadmap,
            screenshots_available=len(screenshots) > 0,
            missing_artifacts=missing_artifacts,
        )

        # Step 4: Create the PPTX file
        await self.log("Building PowerPoint presentation...")
        pptx_bytes = self._build_pptx(slide_content, screenshots)

        # Step 5: Write PPTX to workspace
        await self.write_artifact(PPTX_FILENAME, pptx_bytes)
        artifacts_produced.append(PPTX_FILENAME)
        await self.log(f"Presentation saved: {PPTX_FILENAME}")

        # Step 6: Write speaker notes companion markdown
        speaker_notes_md = self._build_speaker_notes_markdown(slide_content)
        await self.write_artifact(SPEAKER_NOTES_FILENAME, speaker_notes_md)
        artifacts_produced.append(SPEAKER_NOTES_FILENAME)
        await self.log(f"Speaker notes saved: {SPEAKER_NOTES_FILENAME}")

        elapsed = time.monotonic() - start_time
        return AgentResult(
            agent_name=self.agent_name,
            success=True,
            artifacts_produced=artifacts_produced,
            duration_seconds=elapsed,
        )

    async def _read_source_artifact(
        self,
        path: str,
        missing_list: list[str],
        required: bool = True,
    ) -> str:
        """Attempt to read a source artifact from the workspace.

        Args:
            path: Relative path in the workspace.
            missing_list: List to append to if artifact is missing.
            required: Whether this is a required artifact.

        Returns:
            File content or empty string if unavailable.
        """
        try:
            if await self.workspace.file_exists(path):
                content = await self.workspace.read_file(path)
                if content.strip():
                    return content
            if required:
                missing_list.append(path)
                await self.log(f"Source artifact missing: {path}")
            return ""
        except Exception as e:
            if required:
                missing_list.append(path)
                await self.log(f"Error reading {path}: {e}")
            return ""

    async def _find_screenshots(self) -> list[str]:
        """Find any demo screenshot image files in the workspace.

        Looks for image files that may have been captured during QA testing.

        Returns:
            List of relative paths to screenshot files.
        """
        screenshots: list[str] = []
        image_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp"}

        # Check common screenshot locations
        search_dirs = ["screenshots", "video", "qa_screenshots", "logs"]
        for directory in search_dirs:
            try:
                if await self.workspace.file_exists(directory):
                    files = await self.workspace.list_files(directory)
                    for f in files:
                        if Path(f).suffix.lower() in image_extensions:
                            screenshots.append(f)
            except (FileNotFoundError, NotADirectoryError):
                continue

        if screenshots:
            await self.log(f"Found {len(screenshots)} screenshot(s) for slides.")
        else:
            await self.log("No screenshots found. Will use placeholder slides.")

        return screenshots

    async def _generate_slide_content(
        self,
        project_spec: str,
        architecture: str,
        judge_analysis: str,
        roadmap: str,
        screenshots_available: bool,
        missing_artifacts: list[str],
    ) -> list[dict]:
        """Use LLM to generate structured slide content.

        Args:
            project_spec: Contents of project_spec.md.
            architecture: Contents of architecture.md.
            judge_analysis: Contents of judge_analysis.md (may be empty).
            roadmap: Contents of roadmap.md (may be empty).
            screenshots_available: Whether screenshots exist.
            missing_artifacts: List of missing required artifact paths.

        Returns:
            List of slide dictionaries with keys: topic, title, bullets, speaker_notes.
        """
        prompt = self._build_generation_prompt(
            project_spec=project_spec,
            architecture=architecture,
            judge_analysis=judge_analysis,
            roadmap=roadmap,
            screenshots_available=screenshots_available,
            missing_artifacts=missing_artifacts,
        )

        response = await self.llm_generate(prompt, system=SYSTEM_PROMPT)

        # Parse the LLM response as JSON
        slides = self._parse_slide_response(response)

        # Validate and fix up slides
        slides = self._validate_slides(slides, missing_artifacts)

        return slides

    def _build_generation_prompt(
        self,
        project_spec: str,
        architecture: str,
        judge_analysis: str,
        roadmap: str,
        screenshots_available: bool,
        missing_artifacts: list[str],
    ) -> str:
        """Build the LLM prompt for slide content generation.

        Returns:
            Formatted prompt string requesting JSON output.
        """
        parts = [
            "Generate content for a hackathon pitch presentation.",
            "",
            "You MUST return a JSON array of slide objects. Each slide object has:",
            '  - "topic": one of: "title", "problem_statement", "solution_overview", '
            '"technical_architecture", "demo_screenshots", "team_tooling", "future_roadmap"',
            '  - "title": slide title (short, impactful)',
            '  - "bullets": array of 3-5 bullet point strings',
            '  - "speaker_notes": string with talking points (MUST be at least 50 characters)',
            "",
            "REQUIREMENTS:",
            f"- Generate between {MIN_SLIDES} and {MAX_SLIDES} slides total",
            "- You MUST include at least one slide for EACH of these topics:",
            "  problem_statement, solution_overview, technical_architecture,",
            "  demo_screenshots, team_tooling, future_roadmap",
            "- You MAY include a title slide and additional detail slides",
            "- Each speaker_notes field MUST be at least 50 characters long",
            "- Make content compelling, concise, and hackathon-appropriate",
            "",
        ]

        if project_spec:
            parts.extend([
                "--- PROJECT SPECIFICATION ---",
                project_spec[:3000],  # Truncate to fit context
                "",
            ])
        else:
            parts.append(
                "NOTE: project_spec.md is unavailable. "
                "Use generic placeholder content for problem/solution slides.\n"
            )

        if architecture:
            parts.extend([
                "--- ARCHITECTURE ---",
                architecture[:3000],
                "",
            ])
        else:
            parts.append(
                "NOTE: architecture.md is unavailable. "
                "Use generic placeholder content for architecture slide.\n"
            )

        if judge_analysis:
            parts.extend([
                "--- JUDGE ANALYSIS (for context) ---",
                judge_analysis[:1500],
                "",
            ])

        if roadmap:
            parts.extend([
                "--- ROADMAP (for future roadmap slide) ---",
                roadmap[:1500],
                "",
            ])

        if not screenshots_available:
            parts.append(
                "NOTE: No demo screenshots are available. "
                "For the demo_screenshots slide, describe what the demo would show "
                "and note that screenshots will be added when available.\n"
            )

        parts.extend([
            "",
            "Return ONLY the JSON array. No additional text or markdown formatting.",
            "Example format:",
            '[{"topic": "problem_statement", "title": "The Problem", '
            '"bullets": ["Point 1", "Point 2", "Point 3"], '
            '"speaker_notes": "Talk about the core problem we are solving and why it matters to users."}]',
        ])

        return "\n".join(parts)

    def _parse_slide_response(self, response: str) -> list[dict]:
        """Parse the LLM response into a list of slide dictionaries.

        Attempts JSON parsing, with fallback to extracting JSON from markdown.

        Args:
            response: Raw LLM response text.

        Returns:
            List of slide dictionaries.
        """
        # Try direct JSON parse
        try:
            slides = json.loads(response)
            if isinstance(slides, list):
                return slides
        except json.JSONDecodeError:
            pass

        # Try extracting JSON from markdown code blocks
        import re

        json_match = re.search(r"```(?:json)?\s*\n?(.*?)\n?```", response, re.DOTALL)
        if json_match:
            try:
                slides = json.loads(json_match.group(1))
                if isinstance(slides, list):
                    return slides
            except json.JSONDecodeError:
                pass

        # Try finding array brackets in the response
        start = response.find("[")
        end = response.rfind("]")
        if start != -1 and end != -1 and end > start:
            try:
                slides = json.loads(response[start : end + 1])
                if isinstance(slides, list):
                    return slides
            except json.JSONDecodeError:
                pass

        # Fallback: return empty list (will trigger placeholder generation)
        logger.warning("Failed to parse LLM slide response. Using fallback content.")
        return []

    def _validate_slides(
        self, slides: list[dict], missing_artifacts: list[str]
    ) -> list[dict]:
        """Validate and fix slide content to meet requirements.

        Ensures:
        - All required topics are covered
        - Slide count is within bounds (6-15)
        - Speaker notes are at least 50 characters

        Args:
            slides: Parsed slide list from LLM.
            missing_artifacts: Artifacts that were unavailable.

        Returns:
            Validated and potentially augmented slide list.
        """
        # Check which required topics are covered
        covered_topics = set()
        for slide in slides:
            topic = slide.get("topic", "")
            if topic in REQUIRED_TOPICS:
                covered_topics.add(topic)

        # Add placeholder slides for missing topics
        for topic in REQUIRED_TOPICS:
            if topic not in covered_topics:
                placeholder = self._create_placeholder_slide(topic, missing_artifacts)
                slides.append(placeholder)

        # Ensure we have a title slide at the beginning
        has_title = any(s.get("topic") == "title" for s in slides)
        if not has_title:
            slides.insert(0, {
                "topic": "title",
                "title": "Hackathon Project Presentation",
                "bullets": [
                    "Built with AI-powered development",
                    "Autonomous multi-agent collaboration",
                    "From idea to MVP",
                ],
                "speaker_notes": (
                    "Welcome everyone. Today we are presenting our hackathon project "
                    "that was built using an autonomous AI development system."
                ),
            })

        # Enforce slide count bounds
        if len(slides) > MAX_SLIDES:
            # Keep title, required topics, and trim extras
            essential = [s for s in slides if s.get("topic") in REQUIRED_TOPICS or s.get("topic") == "title"]
            extras = [s for s in slides if s not in essential]
            slides = essential + extras[: MAX_SLIDES - len(essential)]

        # Ensure speaker notes meet minimum length
        for slide in slides:
            notes = slide.get("speaker_notes", "")
            if len(notes) < MIN_SPEAKER_NOTES_LENGTH:
                title = slide.get("title", "this slide")
                slide["speaker_notes"] = (
                    f"Key talking points for {title}: "
                    f"Explain the main concepts and highlight the important details "
                    f"that the audience needs to understand about this topic."
                )

        # Ensure all slides have required fields
        for slide in slides:
            if "title" not in slide:
                slide["title"] = "Untitled Slide"
            if "bullets" not in slide or not isinstance(slide["bullets"], list):
                slide["bullets"] = ["Content to be added"]
            if "speaker_notes" not in slide:
                slide["speaker_notes"] = (
                    "Discuss the key points on this slide and engage the audience "
                    "with relevant details about the project."
                )

        return slides

    def _create_placeholder_slide(
        self, topic: str, missing_artifacts: list[str]
    ) -> dict:
        """Create a placeholder slide for a missing topic.

        Args:
            topic: The required topic that needs a placeholder.
            missing_artifacts: List of artifacts that were unavailable.

        Returns:
            Slide dictionary with placeholder content.
        """
        placeholders = {
            "problem_statement": {
                "topic": "problem_statement",
                "title": "The Problem",
                "bullets": [
                    "Problem statement will be detailed here",
                    "Target users face specific challenges",
                    "Current solutions are inadequate",
                    "Opportunity for innovation exists",
                ],
                "speaker_notes": (
                    "Describe the core problem our project addresses. "
                    "Explain who is affected and why existing solutions fall short. "
                    "This sets up the motivation for our solution."
                ),
            },
            "solution_overview": {
                "topic": "solution_overview",
                "title": "Our Solution",
                "bullets": [
                    "Innovative approach to the problem",
                    "Key features and capabilities",
                    "User-centric design philosophy",
                    "Measurable impact and benefits",
                ],
                "speaker_notes": (
                    "Present our solution and its core value proposition. "
                    "Highlight how it addresses each pain point identified in the problem. "
                    "Focus on what makes our approach unique."
                ),
            },
            "technical_architecture": {
                "topic": "technical_architecture",
                "title": "Technical Architecture",
                "bullets": [
                    "System architecture and components",
                    "Technology stack selection",
                    "Integration points and data flow",
                    "Scalability considerations",
                ],
                "speaker_notes": (
                    "Walk through the technical architecture. "
                    "Explain the key technology choices and how components interact. "
                    "Mention scalability and reliability aspects."
                ),
            },
            "demo_screenshots": {
                "topic": "demo_screenshots",
                "title": "Demo & Screenshots",
                "bullets": [
                    "Live demo walkthrough",
                    "Key user interactions",
                    "Screenshots will be added when available",
                    "See the working application in action",
                ],
                "speaker_notes": (
                    "This slide would normally show demo screenshots from QA testing. "
                    "Screenshots were not available during presentation generation. "
                    "Walk through the demo flow verbally or show a live demo instead."
                ),
            },
            "team_tooling": {
                "topic": "team_tooling",
                "title": "Team & Tooling",
                "bullets": [
                    "Built with Hackathon Studio AI system",
                    "10 specialized AI agents collaborating",
                    "Automated testing and quality assurance",
                    "From idea to MVP in hours, not days",
                ],
                "speaker_notes": (
                    "Explain the development process and tools used. "
                    "Highlight the multi-agent AI collaboration that built this project. "
                    "Mention the automated QA and documentation pipeline."
                ),
            },
            "future_roadmap": {
                "topic": "future_roadmap",
                "title": "Future Roadmap",
                "bullets": [
                    "Phase 1: Core feature enhancements",
                    "Phase 2: Scale and performance",
                    "Phase 3: Advanced capabilities",
                    "Long-term vision and growth",
                ],
                "speaker_notes": (
                    "Outline the future plans for the project beyond the hackathon. "
                    "Show a phased roadmap that demonstrates long-term viability. "
                    "Mention stretch goals that could be tackled next."
                ),
            },
        }

        slide = placeholders.get(topic, {
            "topic": topic,
            "title": topic.replace("_", " ").title(),
            "bullets": ["Content placeholder - details to be added"],
            "speaker_notes": (
                "This is a placeholder slide. The content will be filled in "
                "once the relevant source materials become available."
            ),
        })

        # Note if content is unavailable due to missing artifacts
        if missing_artifacts and topic in ("problem_statement", "solution_overview"):
            if "project_spec.md" in missing_artifacts:
                slide["bullets"].append(
                    "[Content unavailable: project_spec.md not found]"
                )
        if missing_artifacts and topic == "technical_architecture":
            if "architecture.md" in missing_artifacts:
                slide["bullets"].append(
                    "[Content unavailable: architecture.md not found]"
                )

        return slide

    def _build_pptx(
        self, slides: list[dict], screenshots: list[str]
    ) -> bytes:
        """Build a PowerPoint file from slide content.

        Args:
            slides: List of slide dictionaries with title, bullets, speaker_notes.
            screenshots: List of screenshot file paths in workspace.

        Returns:
            Bytes of the .pptx file.
        """
        prs = Presentation()

        # Set slide dimensions (widescreen 16:9)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides:
            topic = slide_data.get("topic", "")
            title = slide_data.get("title", "Untitled")
            bullets = slide_data.get("bullets", [])
            speaker_notes = slide_data.get("speaker_notes", "")

            if topic == "title":
                # Use title slide layout
                slide_layout = prs.slide_layouts[0]  # Title Slide
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = title
                # Add subtitle with first bullet
                if len(slide.placeholders) > 1 and bullets:
                    slide.placeholders[1].text = "\n".join(bullets[:2])
            else:
                # Use title and content layout
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                if slide.shapes.title:
                    slide.shapes.title.text = title

                # Add bullet points to the content placeholder
                if len(slide.placeholders) > 1:
                    text_frame = slide.placeholders[1].text_frame
                    text_frame.clear()
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            text_frame.paragraphs[0].text = bullet
                            text_frame.paragraphs[0].font.size = Pt(18)
                        else:
                            p = text_frame.add_paragraph()
                            p.text = bullet
                            p.font.size = Pt(18)
                            p.level = 0

            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = speaker_notes

        # Serialize to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _build_speaker_notes_markdown(self, slides: list[dict]) -> str:
        """Build a markdown companion file with all speaker notes.

        Args:
            slides: List of slide dictionaries.

        Returns:
            Markdown formatted string with speaker notes per slide.
        """
        lines = [
            "# Speaker Notes",
            "",
            "Companion notes for the hackathon presentation slides.",
            "",
        ]

        for i, slide in enumerate(slides, 1):
            title = slide.get("title", "Untitled")
            notes = slide.get("speaker_notes", "")
            lines.extend([
                f"## Slide {i}: {title}",
                "",
                notes,
                "",
            ])

        return "\n".join(lines)
